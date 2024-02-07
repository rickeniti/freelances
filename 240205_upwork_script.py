from pandas import read_csv, DataFrame, ExcelWriter, read_excel
import matplotlib.pyplot as plt
from os import getcwd, path, makedirs
from datetime import datetime
from shutil import copyfile
from openpyxl import Workbook, load_workbook
from pptx import Presentation


### Declaring functions
def replace_chart_in_ppt(presentation_path, slide_index, chart_indices, *chart_image_paths):
    presentation = Presentation(presentation_path)
    slide = presentation.slides[slide_index - 1]

    for chart_index, chart_image_path in zip(chart_indices, chart_image_paths):
        if 0 <= chart_index < len(slide.shapes):
            shape_to_replace = slide.shapes[chart_index]
            shape_to_replace.element.getparent().remove(shape_to_replace.element)
            left, top, width, height = shape_to_replace.left, shape_to_replace.top, shape_to_replace.width, shape_to_replace.height
            slide.shapes.add_picture(chart_image_path, left, top, width, height)

    output_presentation_path = f'{new_folder_path}/{company_name.lower()}_report.pptx'
    presentation.save(output_presentation_path)

def wellness_class(score):
    if score < 45: return 'Surviving'
    elif score > 80: return 'Thriving'
    else: return 'Striving'

def get_current_quarter():
    month = datetime.now().month
    quarters = ['Q1', 'Q2', 'Q3', 'Q4']
    return quarters[(month - 1) // 3]

def category_contributions(series, category, group_avg_score):
    category_count = (series == category).sum() # Counting how many people fall in the same category
    total_participants = len(series) # Getting the total number of responses
    return (category_count / total_participants) * group_avg_score # Calculating their contribution


### Getting inputs from user
company_name = input('Company name: ').upper().replace(" ", "_")
file_name = input('File name: ')


### Creating a new folder
subfolder_name = f'{get_current_quarter()}_{datetime.now().year}'
parent_folder_path = path.join(getcwd(), company_name)
new_folder_path = path.join(parent_folder_path, subfolder_name)

if not path.exists(new_folder_path):
    makedirs(new_folder_path)


### Opening csv file
df = read_csv(f'{file_name}.csv', header=0)


### Renaming columns in the dataset
new_column_names = ['completed','optin','company','vulnerable_group','department','job_role','overall_score','overall_actual','breathing_efficiency_score','breathing_efficiency_actual','movement_habits_score','movement_habits_actual','nutrition_hydration_score','nutrition_hydration_actual','sleep_hygiene_score','sleep_hygiene_actual','quality_of_thinking_score','quality_of_thinking_actual','emotional_regulation_score','emotional_regulation_actual','energy_levels_score','energy_levels_actual','social_connection_score','social_connection_actual','q1','q2','q3','q4','q5','q6','q7','q8','q9','q10','q11','q12','q13','q14','q15','q16','q17','q18','q19','comments']
columns_dict = {}
a = 0

for current_column in df.columns:  
        columns_dict[current_column] = new_column_names[a]
        a += 1

df.rename(columns=columns_dict, inplace=True)


### Classifying each person according to the scores
df['overall_class'] = df['overall_score'].apply(lambda x: wellness_class(x))
df['breathing_class'] = df['breathing_efficiency_score'].apply(lambda x: wellness_class(x))
df['movement_class'] = df['movement_habits_score'].apply(lambda x: wellness_class(x))
df['nutri_hyd_class'] = df['nutrition_hydration_score'].apply(lambda x: wellness_class(x))
df['sleep_hyg_class'] = df['sleep_hygiene_score'].apply(lambda x: wellness_class(x))
df['thinking_class'] = df['quality_of_thinking_score'].apply(lambda x: wellness_class(x))
df['emotional_class'] = df['emotional_regulation_score'].apply(lambda x: wellness_class(x))
df['energy_class'] = df['energy_levels_score'].apply(lambda x: wellness_class(x))
df['social_class'] = df['social_connection_score'].apply(lambda x: wellness_class(x))


# Read the original file using pandas
copyfile('survey_data.xlsx', f'{new_folder_path}/{company_name.lower()}_survey_data.xlsx')
survey_df = read_excel('survey_data.xlsx')

with ExcelWriter(f'{new_folder_path}/{company_name.lower()}_survey_data.xlsx', engine='openpyxl', mode='a', if_sheet_exists='replace') as writer:
    survey_df.to_excel(writer, index=False, sheet_name='Data')

categories = ['Surviving', 'Striving', 'Thriving']
scores_table = df.filter(like='score')
class_table = df.iloc[:,44:]
avg_scores = scores_table.mean().to_dict()
avg_scores_list = scores_table.mean().to_list()

summary_table = DataFrame(index=categories)

for i in range(len(scores_table.columns)):
    temp_table = DataFrame(index=categories, columns=[f'{scores_table.columns[i]}'])
    for category in temp_table.index:
        contribution = category_contributions(class_table.iloc[:,i], category, avg_scores[scores_table.columns[i]])
        temp_table.loc[category] = contribution
    summary_table = concat([summary_table, temp_table], axis = 1)

new_column_names = ['Overall Score %','Breathing Efficiency Score','Movement Habits Score','Nutrition & Hydration Score','Sleep Hygiene Score','Quality Of Thinking Score','Emotional Regulation Score','Energy Levels Score','Social Connection Score']
columns_dict = {}
a = 0

for current_column in summary_table.columns:
        columns_dict[current_column] = new_column_names[a]
        a += 1

summary_table.rename(columns=columns_dict, inplace=True)


### Generating first image
fig, ax = plt.subplots(figsize=(16.97 / 2.54, 10 / 2.54))
bars = summary_table.transpose().plot.barh(stacked=True, ax=ax, color=['#D42034', '#F39125', '#66BC45'], linewidth=0, width=0.75)

ax.legend().set_visible(False)
ax.set_xticks([])
ax.spines['top'].set_visible(False)
ax.spines['right'].set_visible(False)
ax.spines['bottom'].set_visible(False)
ax.spines['left'].set_visible(False)
ax.tick_params(axis='y', which='both', left=False)
ax.set_yticklabels(ax.get_yticklabels(), fontname="Helvetica Neue", fontsize=9)
ax.set_facecolor('none')

totals = summary_table.sum(axis=0)
x_offset = 6

for i, total in enumerate(totals):
  ax.text(total + x_offset, i, f'{round(total)}%', ha='right', va='center',
          fontfamily='Helvetica Neue', fontsize=9)


### Saving image 1
plt.tight_layout()
plt.savefig(f'{new_folder_path}/overall_scores_chart.png', format='png', dpi=300, transparent = True)
plt.close()


### Generating second image
primary = list(avg_scores.values())[0]/100
aux = 1 - primary
share = [aux, primary]
labels = ['', 'Wellness Score']
colors = ['#FFFFFF', '#F39125']

fig, ax = plt.subplots()

wedges, texts, autotexts = ax.pie(share, labels=labels, autopct='%.2f%%',
                                   startangle=90, colors=colors)

circle = plt.Circle(xy=(0, 0), radius=0.75, facecolor='#F0F0EE')
ax.add_artist(circle)
ax.set_facecolor('none')

center_label = ax.text(0, 0, f'{int(share[1] * 100)}%', ha='center', va='center', fontsize=36, color='#F39125', fontweight='extra bold', fontname="Tahoma")

for text in texts + autotexts:
    text.set_visible(False)

ax.axis('equal')


### Saving image 2
plt.tight_layout()
plt.savefig(f'{new_folder_path}/company_wellness_score.png', format='png', dpi=300, transparent = True)
plt.close()


# Specify the PowerPoint presentation path and slide details
presentation_path = 'Template.pptx'
slide_index_to_replace = 2  # Replace with the index of the slide containing the charts
chart_indices_to_replace = [10,7]  # Replace with the indices of the charts to replace on the specified slide
chart_image_paths = [f'{new_folder_path}/overall_scores_chart.png', f'{new_folder_path}/company_wellness_score.png']


# Replace the charts in the PowerPoint presentation
replace_chart_in_ppt(presentation_path, slide_index_to_replace, chart_indices_to_replace, *chart_image_paths)

